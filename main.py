import sys
import os
import uuid
import time
import mimetypes
import platform
from pathlib import Path
from threading import Thread
from subprocess import run

from PySide6.QtWidgets import (
    QApplication, QMainWindow, QLabel, QPushButton, QFileDialog,
    QVBoxLayout, QHBoxLayout, QWidget, QMessageBox,
    QListWidget, QSizePolicy, QListWidgetItem, QDialog,
    QTextEdit, QScrollArea, QLineEdit, QFileIconProvider,
    QListView, QInputDialog, QTableWidget, QTableWidgetItem,
    QMenu
)
from PySide6.QtGui     import QFont, QPixmap, QImage, QFontMetrics, QAction
from PySide6.QtCore import Qt, QSize, QFileInfo

from cryptography.hazmat.primitives.ciphers import Cipher, algorithms, modes
from cryptography.hazmat.primitives import padding
from cryptography.hazmat.backends import default_backend

try:
    from PySide6.QtMultimediaWidgets import QCameraViewfinder
    from PySide6.QtMultimedia import QCamera, QImageCapture, QMediaDevices
except ImportError:
    QCameraViewfinder = None

import mammoth
import pandas as pd
from pptx import Presentation

IMPORT_DIR = Path("imported_docs")
IMPORT_DIR.mkdir(exist_ok=True)
TEMP_DIR = IMPORT_DIR / "tmp"
TEMP_DIR.mkdir(exist_ok=True)

# unterstützte Office-Formate für headless → PDF
office = {".doc", ".docx", ".odt", ".rtf",
          ".xls", ".xlsx", ".ods",
          ".ppt", ".pptx", ".odp"}
WATCHED_PATHS_FILE = Path("watched_folders.txt")
def load_watched_folders():
    if WATCHED_PATHS_FILE.exists():
        return [Path(line.strip()) for line in WATCHED_PATHS_FILE.read_text(encoding="utf-8").splitlines() if line.strip()]
    else:
        default = Path("watch_folder")
        default.mkdir(exist_ok=True)
        return [default]

def save_watched_folders(folders):
    WATCHED_PATHS_FILE.write_text('\n'.join(str(f) for f in folders), encoding="utf-8")

KEY = os.environ.get("DMS_AES_KEY", None)
if not KEY:
    KEY = b"0123456789abcdef0123456789abcdef"
else:
    KEY = KEY.encode("utf-8")
IV_SIZE = 16

def encrypt_file(src_path, dest_path, key):
    backend = default_backend()
    iv = os.urandom(IV_SIZE)
    cipher = Cipher(algorithms.AES(key), modes.CBC(iv), backend=backend)
    encryptor = cipher.encryptor()
    padder = padding.PKCS7(128).padder()
    with open(src_path, "rb") as f_in, open(dest_path, "wb") as f_out:
        f_out.write(iv)
        while True:
            chunk = f_in.read(1024 * 64)
            if not chunk:
                break
            padded = padder.update(chunk)
            if padded:
                f_out.write(encryptor.update(padded))
        padded = padder.finalize()
        if padded:
            f_out.write(encryptor.update(padded))
        f_out.write(encryptor.finalize())

def decrypt_file(src_path, dest_path, key):
    backend = default_backend()
    with open(src_path, "rb") as f_in:
        iv = f_in.read(IV_SIZE)
        cipher = Cipher(algorithms.AES(key), modes.CBC(iv), backend=backend)
        decryptor = cipher.decryptor()
        unpadder = padding.PKCS7(128).unpadder()
        with open(dest_path, "wb") as f_out:
            while True:
                chunk = f_in.read(1024 * 64)
                if not chunk:
                    break
                decrypted = decryptor.update(chunk)
                if decrypted:
                    f_out.write(unpadder.update(decrypted))
            decrypted = decryptor.finalize()
            if decrypted:
                f_out.write(unpadder.update(decrypted))
            try:
                f_out.write(unpadder.finalize())
            except ValueError:
                pass  # Falls Padding fehlschlägt

def create_encrypted_temp_file(data: bytes, key: bytes, suffix: str = ".tmp.enc") -> Path:
    """
    Legt eine temporäre, verschlüsselte Datei im TEMP_DIR an und gibt den Pfad zurück.
    Die Datei wird mit AES-256 verschlüsselt.
    """
    temp_filename = f"{uuid.uuid4().hex}{suffix}"
    temp_path = TEMP_DIR / temp_filename
    backend = default_backend()
    iv = os.urandom(IV_SIZE)
    cipher = Cipher(algorithms.AES(key), modes.CBC(iv), backend=backend)
    encryptor = cipher.encryptor()
    padder = padding.PKCS7(128).padder()
    with open(temp_path, "wb") as f_out:
        f_out.write(iv)
        padded = padder.update(data) + padder.finalize()
        f_out.write(encryptor.update(padded) + encryptor.finalize())
    return temp_path

def delete_temp_file(temp_path: Path):
    """
    Löscht eine temporäre Datei sicher aus dem TEMP_DIR.
    """
    try:
        temp_path.unlink()
    except Exception:
        pass

class MainWindow(QMainWindow):
    def __init__(self):
        super().__init__()
        self.setWindowTitle("DMS – Dokumentenmanagementsystem")
        self.resize(1000, 700)

        self.watched_folders = load_watched_folders()
        self.preview_dialog = None

        # alle UI-Elemente und Layouts anlegen
        self.setup_ui()

        # Ordner-Watcher starten
        self.watching = True
        Thread(target=self.watch_folders, daemon=True).start()

    def setup_ui(self):
        # ─── NAVIGATION BUTTONS ─────────────────────────────────────────────
        nav = QVBoxLayout()
        for label, slot in [
            ("Scannen", self.scan_document),
            ("Importieren", self.import_document),
            ("Exportieren", self.export_document),
            ("Vorschau", self.preview_document),
            ("Einstellungen", lambda: None)
        ]:
            btn = QPushButton(label)
            btn.setFont(QFont("Arial", 14))
            btn.clicked.connect(slot)
            nav.addWidget(btn)
            if label == "Vorschau":
                self.btn_preview = btn
                btn.setEnabled(False)
        nav.addStretch()
        nav_widget = QWidget()
        nav_widget.setLayout(nav)
        nav_widget.setFixedWidth(180)

        # ─── ORDNER-NAVIGATION (echte Unterordner + „Alle Dokumente“) ─────────
        # nutze QFileIconProvider, um ein konsistentes Ordner-Icon zu holen
        provider = QFileIconProvider()
        dir_icon = provider.icon(QFileIconProvider.Folder)

        # virtuellen Ordner-Hinzufügen-Button
        folder_box = QVBoxLayout()
        btn_new_folder = QPushButton("+ Ordner")
        btn_new_folder.setFixedWidth(80)
        btn_new_folder.clicked.connect(self.add_virtual_folder)
        folder_box.addWidget(btn_new_folder, 0, Qt.AlignHCenter)

        self.folder_nav = QListWidget()
        self.folder_nav.setIconSize(QSize(16,16))
        # Wurzel-Eintrag
        self.folder_nav.addItem(QListWidgetItem(dir_icon, "Alle Dokumente"))
        # alle Unterordner in IMPORT_DIR
        for d in sorted(IMPORT_DIR.iterdir()):
            if d.is_dir():
                self.folder_nav.addItem(QListWidgetItem(dir_icon, d.name))
        self.folder_nav.currentItemChanged.connect(self.on_folder_selected)
        self.folder_nav.setFixedWidth(180)
        folder_box.addWidget(self.folder_nav)
        folder_widget = QWidget()
        folder_widget.setLayout(folder_box)
        folder_widget.setFixedWidth(200)

        # ─── DATEILISTE ────────────────────────────────────────────────────
        self.doc_list = QListWidget()
        # Drag & Drop und Kontextmenü aktivieren
        self.doc_list.setDragEnabled(True)
        self.doc_list.setAcceptDrops(True)
        self.doc_list.setDefaultDropAction(Qt.MoveAction)
        self.doc_list.setContextMenuPolicy(Qt.CustomContextMenu)
        self.doc_list.customContextMenuRequested.connect(self.show_doc_context_menu)
        self.doc_list.setViewMode(QListWidget.IconMode)
        self.doc_list.setIconSize(QSize(64,64))
        self.doc_list.setResizeMode(QListWidget.Adjust)
        # echtes Grid-Layout mit Zeilenumbruch und Text-Umbruch
        self.doc_list.setWordWrap(True)
        self.doc_list.setFlow(QListView.LeftToRight)
        self.doc_list.setWrapping(True)
        self.doc_list.setUniformItemSizes(True)
        self.doc_list.currentItemChanged.connect(self.on_file_selected)
 
        # ─── DATEI-INFOS + KLEINE PREVIEW ─────────────────────────────────
        info = QVBoxLayout()
        self.info_name = QLabel("Name: –")
        self.info_size = QLabel("Größe: –")
        info.addWidget(self.info_name)
        info.addWidget(self.info_size)
        self.preview_small = QLabel()
        self.preview_small.setFixedSize(120,120)
        self.preview_small.setStyleSheet("border:1px solid #aaa;")
        self.preview_small.mousePressEvent = lambda e: self.show_large_preview()
        info.addWidget(self.preview_small)
        info.addStretch()
        info_widget = QWidget()
        info_widget.setLayout(info)
        info_widget.setFixedWidth(200)

        # ─── WATCH-FOLDER VERWALTUNG ───────────────────────────────────────
        bottom = QHBoxLayout()
        self.folder_list = QListWidget()
        self.refresh_folder_list()
        bottom.addWidget(self.folder_list)
        btn_add = QPushButton("+"); btn_add.clicked.connect(self.add_watch_folder)
        btn_rem = QPushButton("–"); btn_rem.clicked.connect(self.remove_selected_watch_folder)
        bottom.addWidget(btn_add); bottom.addWidget(btn_rem)
        bottom_widget = QWidget()
        bottom_widget.setLayout(bottom)

        # ─── LAYOUT ZUSAMMENSETZEN ────────────────────────────────────────
        top = QHBoxLayout()
        top.addWidget(nav_widget)
        top.addWidget(folder_widget)
        top.addWidget(self.doc_list, 1)    # <— hier das Icon‐Grid einfügen
        top.addWidget(info_widget)         # <— hier die Info‐Spalte einfügen

        main = QVBoxLayout()
        main.addLayout(top, 1)
        main.addWidget(bottom_widget)

        container = QWidget()
        container.setLayout(main)
        self.setCentralWidget(container)

        # initiale Befüllung der Dateiliste (erst jetzt existiert info_name)
        self.refresh_doc_list()

    # ─── SLOT-METHODEN (korrekt eingerückt!) ─────────────────────────────
    def on_folder_selected(self, current, previous):
        if not current:
            return
        text = current.text()
        if text == "Alle Dokumente":
            self.refresh_doc_list()
        else:
            path = IMPORT_DIR / text
            if path.is_dir():
                self.refresh_doc_list(path)

    def preview_document(self):
        item = self.doc_list.currentItem()
        if not item:
            if self.preview_dialog:
                self.preview_dialog.close()
            return

        # Entschlüssele in eine temporäre Datei
        enc_path: Path = item.data(Qt.UserRole)
        # Icon-Provider schon jetzt bereitstellen (für CSV/Default-Fälle)
        provider = QFileIconProvider()
        from tempfile import NamedTemporaryFile
        suffix = Path(enc_path.stem).suffix or ""
        with NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
            tmp_path = Path(tmp.name)
        try:
            decrypt_file(enc_path, tmp_path, KEY)

            # großen Vorschau-Dialog aufbauen
            dlg = QDialog(self)
            # Referenz zum Schließen später behalten
            self.preview_dialog = dlg
            dlg.setWindowTitle(enc_path.stem)
            dlg.resize(800, 600)
            layout = QVBoxLayout(dlg)

            ext = tmp_path.suffix.lower()
            # TXT-, MD-, PY-Dateien → reiner Text
            if ext in (".txt", ".md", ".py"):
                txt = tmp_path.read_text(encoding="utf-8", errors="ignore")
                te = QTextEdit(); te.setReadOnly(True); te.setPlainText(txt)
                layout.addWidget(te); dlg.exec(); return
            # DOCX → HTML per Mammoth
            if ext == ".docx":
                # Temp-File sauber öffnen und schließen
                with tmp_path.open("rb") as f:
                    html = mammoth.convert_to_html(f).value
                te = QTextEdit(); te.setReadOnly(True); te.setHtml(html)
                layout.addWidget(te); dlg.exec(); return
            # XLSX/ODS → HTML via Pandas
            elif ext in (".xls", ".xlsx", ".ods"):
                df = pd.read_excel(tmp_path, engine="openpyxl")
                html = df.to_html(index=False)
                te = QTextEdit(); te.setReadOnly(True); te.setHtml(html)
                layout.addWidget(te); dlg.exec(); return
            # PPTX/ODP → Einfacher HTML-Text via python-pptx
            elif ext in (".ppt", ".pptx", ".odp"):
                prs = Presentation(str(tmp_path))
                html = ""
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            html += "<p>" + shape.text.replace("\n","<br>") + "</p>"
                te = QTextEdit(); te.setReadOnly(True); te.setHtml(html or "<p>Keine Textinhalte.</p>")
                layout.addWidget(te); dlg.exec(); return
            # CSV → Tabelle
            if ext == ".csv":
                import csv
                table = QTableWidget()
                # CSV einlesen
                with open(tmp_path, newline="", encoding="utf-8", errors="ignore") as f:
                    reader = csv.reader(f)
                    data = list(reader)
                if data:
                    table.setColumnCount(len(data[0]))
                    table.setRowCount(len(data))
                    for r, row in enumerate(data):
                        for c, cell in enumerate(row):
                            table.setItem(r, c, QTableWidgetItem(cell))
                    table.resizeColumnsToContents()
                layout.addWidget(table)
                dlg.exec()
                return

            # Pixmap für Grafik, PDF, Office oder Icon ermitteln
            if ext in (".png", ".jpg", ".jpeg", ".bmp", ".gif"):
                pix = QPixmap(str(tmp_path))
            elif ext in office:
                # Office → Headless PDF → erste Seite
                pdf_tmp = tmp_path.with_suffix(".pdf")
                run(["soffice","--headless","--convert-to","pdf",
                     "--outdir", str(pdf_tmp.parent), str(tmp_path)], check=False)
                import fitz
                doc = fitz.open(str(pdf_tmp))
                pm = doc.load_page(0).get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
                img = QImage(pm.samples, pm.width, pm.height, pm.stride, QImage.Format_RGB888)
                pix = QPixmap.fromImage(img)
                doc.close()
            elif ext == ".pdf":
                import fitz
                doc = fitz.open(str(tmp_path))
                pm = doc.load_page(0).get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
                img = QImage(pm.samples, pm.width, pm.height, pm.stride, QImage.Format_RGB888)
                pix = QPixmap.fromImage(img)
                doc.close()
            else:
                # Fallback-Icon
                icon = provider.icon(QFileInfo(tmp_path))
                pix = icon.pixmap(800, 600)

            # alle Grafik-/PDF-Fälle: zoombarer View
            from PySide6.QtWidgets import QGraphicsView, QGraphicsScene, QSlider, QLabel, QHBoxLayout
            # Zoom‐Control
            ctrl = QHBoxLayout()
            ctrl.addWidget(QLabel("Zoom %:"))
            slider = QSlider(Qt.Horizontal)
            slider.setRange(10, 400)
            slider.setValue(100)
            ctrl.addWidget(slider)
            layout.addLayout(ctrl)

            # Grafik-View
            view = QGraphicsView()
            scene = QGraphicsScene(view)
            item = scene.addPixmap(pix)
            view.setScene(scene)
            view.setDragMode(QGraphicsView.ScrollHandDrag)
            view.setTransformationAnchor(QGraphicsView.AnchorUnderMouse)
            layout.addWidget(view, 1)

            # Zoom-Funktion verbinden
            def on_zoom(val):
                view.resetTransform()
                factor = val / 100.0
                view.scale(factor, factor)
            slider.valueChanged.connect(on_zoom)

            dlg.exec()
        finally:
            # versuche zu löschen, ignoriere falls noch in Benutzung
            try:
                tmp_path.unlink()
            except Exception:
                pass

    def show_large_preview(self):
        self.preview_document()

    def remove_selected_watch_folder(self):
        item = self.folder_list.currentItem()
        if not item:
            return
        # hier dein Entfernen-Code …

    # ─── BISHERIGE FUNKTIONEN UNVERÄNDERT EINBAUEN ───────────────────────
    def refresh_doc_list(self, folder: Path = None):
        """
        Zeigt im Icon-Grid alle *.enc-Dateien in IMPORT_DIR (oder Unterordner),
        aber elidiert die Anzeige ohne '.enc' – voller Name als Tooltip.
        """
        folder = folder or IMPORT_DIR
        self.doc_list.clear()
        provider = QFileIconProvider()
        fm = QFontMetrics(self.doc_list.font())
        max_width = 100  # px für den elidierten Text
        for entry in sorted(folder.glob("*.enc")):
            # Originalname und Extension
            orig = entry.stem  # z.B. "Report.docx"
            # passendes Icon nach Original-Extension
            icon = provider.icon(QFileInfo(orig))
            elided = fm.elidedText(orig, Qt.ElideRight, max_width)
            item = QListWidgetItem(icon, elided)
            # speichere das komplette Path-Objekt
            item.setData(Qt.UserRole, entry)
            item.setToolTip(orig)
            self.doc_list.addItem(item)
 
    def on_file_selected(self, current, previous):
        """
        Wenn ein Icon in der Mitte ausgewählt wird, entschlüsseln wir
        und zeigen Name, Größe und kleine Vorschau in der rechten Spalte.
        """
        if not current:
            return
        enc_path: Path = current.data(Qt.UserRole)
        # für kleine Vorschau in TEMP entschlüsseln und laden
        from tempfile import NamedTemporaryFile
        # suffix aus Original-Extension
        suffix = Path(enc_path.stem).suffix or ""
        with NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
            tmp_path = Path(tmp.name)
        try:
            decrypt_file(enc_path, tmp_path, KEY)
            # Dateiendung für Vorschau ermitteln
            ext = tmp_path.suffix.lower()
            provider = QFileIconProvider()
            if ext in (".png", ".jpg", ".jpeg", ".bmp", ".gif"):
                pix = QPixmap(str(tmp_path))
            # PDF-Formate: erste Seite als Bild
            elif ext == ".pdf":
                import fitz
                doc = fitz.open(str(tmp_path))
                # schärfer rendern mit höherem Zoom
                mat = fitz.Matrix(2, 2)
                pm = doc.load_page(0).get_pixmap(matrix=mat, alpha=False)
                img = QImage(pm.samples, pm.width, pm.height, pm.stride, QImage.Format_RGB888)
                pix = QPixmap.fromImage(img)
                doc.close()
            # CSV bekommt ein generisches Icon
            elif ext == ".csv":
                icon = provider.icon(QFileInfo(tmp_path))
                pix = icon.pixmap(120, 120)
            # alle anderen: Standard-Icon
            else:
                icon = provider.icon(QFileInfo(tmp_path))
                pix = icon.pixmap(120, 120)

            # Name & Größe anzeigen
            display = enc_path.stem
            self.info_name.setText(f"Name: {display}")
            size = tmp_path.stat().st_size
            self.info_size.setText(f"Größe: {size} Bytes")
            if not pix.isNull():
                self.preview_small.setPixmap(pix.scaled(120, 120, Qt.KeepAspectRatio, Qt.SmoothTransformation))
            else:
                self.preview_small.clear()
        finally:
            try:
                tmp_path.unlink()
            except:
                pass

    def show_large_preview(self):
        self.preview_document()

    def remove_selected_watch_folder(self):
        item = self.folder_list.currentItem()
        if not item:
            return
        # hier dein Entfernen-Code …

    # ─── BISHERIGE FUNKTIONEN UNVERÄNDERT EINBAUEN ───────────────────────
    def refresh_doc_list(self, folder: Path = None):
        """
        Zeigt im Icon-Grid alle *.enc-Dateien in IMPORT_DIR (oder Unterordner),
        aber elidiert die Anzeige ohne '.enc' – voller Name als Tooltip.
        """
        folder = folder or IMPORT_DIR
        self.doc_list.clear()
        provider = QFileIconProvider()
        fm = QFontMetrics(self.doc_list.font())
        max_width = 100  # px für den elidierten Text
        for entry in sorted(folder.glob("*.enc")):
            # Originalname und Extension
            orig = entry.stem  # z.B. "Report.docx"
            # passendes Icon nach Original-Extension
            icon = provider.icon(QFileInfo(orig))
            elided = fm.elidedText(orig, Qt.ElideRight, max_width)
            item = QListWidgetItem(icon, elided)
            # speichere das komplette Path-Objekt
            item.setData(Qt.UserRole, entry)
            item.setToolTip(orig)
            self.doc_list.addItem(item)
 
    def refresh_folder_list(self):
        self.folder_list.clear()
        for folder in self.watched_folders:
            item = QListWidgetItem(str(folder.resolve()))
            self.folder_list.addItem(item)
        # Kontextmenü für Entfernen
        self.folder_list.setContextMenuPolicy(Qt.CustomContextMenu)
        self.folder_list.customContextMenuRequested.connect(self.show_folder_context_menu)

    def show_folder_context_menu(self, pos):
        item = self.folder_list.itemAt(pos)
        if item is None:
            return
        from PySide6.QtWidgets import QMenu
        menu = QMenu(self.folder_list)
        remove_action = menu.addAction("Ordner aus Überwachung entfernen")
        action = menu.exec(self.folder_list.mapToGlobal(pos))
        if action == remove_action:
            idx = self.folder_list.row(item)
            if 0 <= idx < len(self.watched_folders):
                removed = self.watched_folders.pop(idx)
                save_watched_folders(self.watched_folders)
                self.refresh_folder_list()

    def add_watch_folder(self):
        dialog = QFileDialog(self)
        dialog.setFileMode(QFileDialog.Directory)
        dialog.setOption(QFileDialog.ShowDirsOnly, True)
        if dialog.exec():
            selected = dialog.selectedFiles()
            if selected:
                folder = Path(selected[0])
                if folder not in self.watched_folders:
                    self.watched_folders.append(folder)
                    save_watched_folders(self.watched_folders)
                    self.refresh_folder_list()

    def import_document(self):
        file_dialog = QFileDialog(self)
        file_dialog.setFileMode(QFileDialog.ExistingFile)
        if file_dialog.exec():
            selected_files = file_dialog.selectedFiles()
            if selected_files:
                src_path = Path(selected_files[0])
                base_name = src_path.name
                dest_path = IMPORT_DIR / (base_name + ".enc")
                # Prüfe, ob schon eine Version existiert
                if dest_path.exists():
                    # Ermittle höchste Version
                    stem, ext = os.path.splitext(base_name)
                    versions = []
                    for file in IMPORT_DIR.glob(f"{stem}.v*.{ext}.enc"):
                        pass  # ... wie gehabt ...
                    if versions:
                        max_v = max([int(p.name.split(".v")[-1].split(f".{ext}.enc")[0]) for p in versions])
                    else:
                        max_v = 1
                    # Alte Version umbenennen
                    dest_path.rename(IMPORT_DIR / f"{stem}.v{max_v+1}.{ext}.enc")
                encrypt_file(src_path, dest_path, KEY)
                QMessageBox.information(self, "Importiert", f"Datei wurde importiert und verschlüsselt gespeichert:\n{dest_path}")
                self.refresh_doc_list()

    def export_document(self):
        selected = self.doc_list.currentItem()
        if not selected:
            QMessageBox.warning(self, "Kein Dokument ausgewählt", "Bitte wählen Sie ein Dokument aus der Liste aus.")
            return
        src_path = IMPORT_DIR / selected.text()
        # Zielpfad wählen
        save_dialog = QFileDialog(self)
        save_dialog.setAcceptMode(QFileDialog.AcceptSave)
        save_dialog.setDirectory(str(Path.home()))
        orig_name = src_path.stem
        save_dialog.selectFile(orig_name)
        if save_dialog.exec():
            dest_files = save_dialog.selectedFiles()
            if dest_files:
                dest_path = Path(dest_files[0])
                try:
                    decrypt_file(src_path, dest_path, KEY)
                    QMessageBox.information(self, "Exportiert", f"Datei wurde entschlüsselt exportiert:\n{dest_path}")
                except Exception as e:
                    QMessageBox.critical(self, "Fehler", f"Fehler beim Exportieren:\n{e}")

    def enable_preview_button(self):
        item = self.doc_list.currentItem()
        self.btn_preview.setEnabled(item is not None)

    

    def delete_document(self, item):
        path: Path = item.data(Qt.UserRole)
        if QMessageBox.question(self, "Löschen", f"Datei wirklich löschen?\n{path.name}") == QMessageBox.Yes:
            try:
                path.unlink()
                self.refresh_doc_list()
                QMessageBox.information(self, "Gelöscht", f"{path.name} wurde gelöscht.")
            except Exception as e:
                QMessageBox.critical(self, "Fehler", f"Löschen fehlgeschlagen:\n{e}")

    def move_document(self, item, target_folder: str):
        src: Path = item.data(Qt.UserRole)
        if target_folder == "Alle Dokumente":
            dst_dir = IMPORT_DIR
        else:
            dst_dir = IMPORT_DIR / target_folder
        dst_dir.mkdir(exist_ok=True)
        dst = dst_dir / src.name
        try:
            src.replace(dst)
            self.refresh_doc_list()
            QMessageBox.information(self, "Verschoben", f"{src.name} → {target_folder}")
        except Exception as e:
            QMessageBox.critical(self, "Fehler", f"Verschieben fehlgeschlagen:\n{e}")

    def show_doc_context_menu(self, pos):
        """Rechtsklick-Kontextmenü für die Dateiliste."""
        item = self.doc_list.itemAt(pos)
        if not item:
            return
        menu = QMenu(self.doc_list)
        delete_act = QAction("Löschen", self)
        move_act   = QAction("Verschieben…", self)
        menu.addAction(delete_act)
        menu.addAction(move_act)
        action = menu.exec(self.doc_list.mapToGlobal(pos))
        if action == delete_act:
            self.delete_document(item)
        elif action == move_act:
            folders = ["Alle Dokumente"] + [d.name for d in IMPORT_DIR.iterdir() if d.is_dir()]
            tgt, ok = QInputDialog.getItem(self, "Verschieben", "Zielordner:", folders, 0, False)
            if ok:
                self.move_document(item, tgt)

    def watch_folders(self):
        already_seen = {}
        while self.watching:
            for folder in self.watched_folders:
                if folder not in already_seen:
                    try:
                        already_seen[folder] = set(p.name for p in folder.iterdir() if p.is_file())
                    except Exception:
                        already_seen[folder] = set()
            for folder in list(already_seen.keys()):
                if folder not in self.watched_folders:
                    del already_seen[folder]
            for folder in self.watched_folders:
                if not folder.exists():
                    continue
                try:
                    current_files = set(p.name for p in folder.iterdir() if p.is_file())
                except Exception:
                    continue
                new_files = current_files - already_seen.get(folder, set())
                for filename in new_files:
                    src_path = folder / filename
                    dest_path = IMPORT_DIR / (filename + ".enc")
                    try:
                        encrypt_file(src_path, dest_path, KEY)
                        src_path.unlink()  # Nach Import löschen
                        self.refresh_doc_list()
                    except Exception as e:
                        print(f"Fehler beim automatischen Import aus {folder}: {e}")
                already_seen[folder] = current_files
            time.sleep(2)

    def closeEvent(self, event):
        self.watching = False
        super().closeEvent(event)

    def scan_document(self):
        os_type = platform.system().lower()
        if os_type == "windows":
            try:
                import win32com.client
                wia = win32com.client.Dispatch("WIA.CommonDialog")
                device = wia.ShowSelectDevice()
                if device is None:
                    QMessageBox.information(self, "Scanner", "Kein Scanner ausgewählt.")
                    return
                item = device.Items[1]
                image = wia.ShowAcquireImage(device)
                if image:
                    # Speichere das Bild temporär
                    temp_path = Path("scanned_image.jpg")
                    image.SaveFile(str(temp_path))
                    QMessageBox.information(self, "Scanner", f"Scan gespeichert: {temp_path}")
                    # TODO: Importiere das Bild wie ein gescanntes Dokument
                else:
                    QMessageBox.warning(self, "Scanner", "Scan fehlgeschlagen.")
            except Exception as e:
                QMessageBox.critical(self, "Scanner", f"Fehler beim Scannen: {e}")
        else:
            QMessageBox.information(self, "Scanner", "Scanner-Integration ist nur unter Windows (WIA) vorbereitet.\nFür Linux/Mac bitte SANE/TWAIN-Anbindung ergänzen.")

    def search_documents(self, text):
        import pytesseract
        from PIL import Image
        text = text.lower()
        self.doc_list.clear()
        print("DEBUG: Suche gestartet")
        if not text.strip():
            self.refresh_doc_list()
            return
        for file in sorted(IMPORT_DIR.glob("*.enc")):
            print(f"DEBUG: Datei gefunden: {file.name}")
            if text in file.name.lower():
                self.doc_list.addItem(file.name)
                continue
            orig_name = file.name
            if orig_name.endswith('.enc'):
                orig_name = orig_name[:-4]
            ext = Path(orig_name).suffix.lower()
            from tempfile import NamedTemporaryFile
            with NamedTemporaryFile(delete=False, suffix=ext) as tmp:
                temp_path = Path(tmp.name)
            try:
                decrypt_file(file, temp_path, KEY)
                if ext in [".txt", ".md", ".py"]:
                    with open(temp_path, "r", encoding="utf-8", errors="ignore") as f:
                        content = f.read().lower()
                    if text in content:
                        self.doc_list.addItem(file.name)
                        continue
                elif ext == ".csv":
                    import csv
                    with open(temp_path, "r", encoding="utf-8", errors="ignore") as f:
                        sample = f.read(2048)
                        f.seek(0)
                        try:
                            dialect = csv.Sniffer().sniff(sample)
                        except Exception:
                            dialect = csv.excel  # Fallback
                        reader = csv.reader(f, dialect)
                        for row in reader:
                            if any(text in (cell or "").lower() for cell in row):
                                self.doc_list.addItem(file.name)
                                break
                elif ext == ".pdf":
                    try:
                        import fitz
                        doc = fitz.open(temp_path)
                        pdf_text = ""
                        for page in doc:
                            page_text = page.get_text()
                            print(f"DEBUG: PDF Seite Text: {page_text[:200]}")
                            pdf_text += page_text.lower()
                        doc.close()
                        print(f"DEBUG: Gesamter PDF-Text: {pdf_text[:500]}")
                        print(f"DEBUG: Suche nach '{text}' in PDF-Text: {text in pdf_text}")
                        if text in pdf_text:
                            print(f"DEBUG: Treffer im PDF-Text ({file.name})")
                            self.doc_list.addItem(file.name)
                            continue
                        # OCR für PDF-Seiten (nur wenn kein Text gefunden)
                        doc = fitz.open(temp_path)
                        for page in doc:
                            pix = page.get_pixmap()
                            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                            ocr_text = pytesseract.image_to_string(img, lang="deu+eng").lower()
                            print(f"DEBUG: OCR-Text: {ocr_text[:200]}")
                            print(f"DEBUG: Suche nach '{text}' in OCR-Text: {text in ocr_text}")
                            if text in ocr_text:
                                print(f"DEBUG: Treffer im OCR ({file.name})")
                                self.doc_list.addItem(file.name)
                                break
                        doc.close()
                    except Exception as e:
                        print(f"DEBUG: Fehler bei PDF-Suche: {e}")
                elif ext in [".jpg", ".jpeg", ".png", ".bmp", ".gif"]:
                    try:
                        img = Image.open(temp_path)
                        ocr_text = pytesseract.image_to_string(img, lang="deu+eng").lower()
                        if text in ocr_text:
                            self.doc_list.addItem(file.name)
                    except Exception:
                        pass
            except Exception as e:
                print(f"DEBUG: Fehler bei Datei {file.name}: {e}")
            finally:
                try:
                    temp_path.unlink(missing_ok=True)
                except Exception:
                    pass

    def add_virtual_folder(self):
        """Legt einen neuen Unterordner in imported_docs an und aktualisiert die Nav."""
        name, ok = QInputDialog.getText(self, "Neuer Ordner", "Name des neuen Ordners:")
        if not (ok and name.strip()):
            return
        new_path = IMPORT_DIR / name.strip()
        try:
            new_path.mkdir(exist_ok=False)
        except Exception as e:
            QMessageBox.critical(self, "Fehler", f"Ordner konnte nicht erstellt werden:\n{e}")
            return
        # Neu in die Nav einfügen und sichtbar machen
        provider = QFileIconProvider()
        icon = provider.icon(QFileIconProvider.Folder)
        self.folder_nav.addItem(QListWidgetItem(icon, name.strip()))


if __name__ == "__main__":
    app = QApplication(sys.argv)
    window = MainWindow()
    window.show()
    sys.exit(app.exec())